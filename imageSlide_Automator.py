import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
import os

class SlideMateApp:
    def __init__(self, root):
        self.root = root
        self.root.title("SlideMate")

        # フォルダパスとパワポパスを表示するエントリーボックス
        self.folder_path = tk.Entry(root, width=50)
        self.folder_path.grid(row=0, column=1, padx=10, pady=5)
        self.ppt_path = tk.Entry(root, width=50)
        self.ppt_path.grid(row=1, column=1, padx=10, pady=5)

        # フォルダ選択ボタン
        self.folder_button = tk.Button(root, text="画像フォルダを選択", command=self.select_folder)
        self.folder_button.grid(row=0, column=0, padx=10, pady=5)

        # パワポ選択ボタン
        self.ppt_button = tk.Button(root, text="PowerPointファイルを選択", command=self.select_ppt_file)
        self.ppt_button.grid(row=1, column=0, padx=10, pady=5)

        # 画像サイズの設定
        tk.Label(root, text="画像の幅 (cm):").grid(row=2, column=0, padx=10, pady=5)
        self.img_width = tk.Entry(root, width=10)
        self.img_width.grid(row=2, column=1, padx=10, pady=5)
        self.img_width.insert(0, "12.7")  # デフォルト値（5インチ ≈ 12.7cm）

        tk.Label(root, text="画像の高さ (cm):").grid(row=3, column=0, padx=10, pady=5)
        self.img_height = tk.Entry(root, width=10)
        self.img_height.grid(row=3, column=1, padx=10, pady=5)
        self.img_height.insert(0, "12.7")  # デフォルト値

        # 画像位置の設定
        tk.Label(root, text="画像の左位置 (cm):").grid(row=4, column=0, padx=10, pady=5)
        self.img_left = tk.Entry(root, width=10)
        self.img_left.grid(row=4, column=1, padx=10, pady=5)
        self.img_left.insert(0, "2.54")  # デフォルト値（1インチ ≈ 2.54cm）

        tk.Label(root, text="画像の上位置 (cm):").grid(row=5, column=0, padx=10, pady=5)
        self.img_top = tk.Entry(root, width=10)
        self.img_top.grid(row=5, column=1, padx=10, pady=5)
        self.img_top.insert(0, "2.54")  # デフォルト値

        # 実行ボタン
        self.run_button = tk.Button(root, text="実行", command=self.insert_images_to_ppt)
        self.run_button.grid(row=6, column=0, columnspan=2, padx=10, pady=10)

    def select_folder(self):
        folder_selected = filedialog.askdirectory()
        if folder_selected:
            self.folder_path.delete(0, tk.END)
            self.folder_path.insert(0, folder_selected)

    def select_ppt_file(self):
        file_selected = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
        if file_selected:
            self.ppt_path.delete(0, tk.END)
            self.ppt_path.insert(0, file_selected)

    def clone_slide(self, source_slide, presentation):
        """
        スライドを複製する関数
        """
        # 元のスライドのレイアウトを取得
        slide_layout = source_slide.slide_layout

        # 新しいスライドを追加
        new_slide = presentation.slides.add_slide(slide_layout)

        # すべてのシェイプをコピー
        for shape in source_slide.shapes:
            if shape.has_text_frame:
                # テキストボックスを複製
                new_shape = new_slide.shapes.add_shape(
                    shape.auto_shape_type,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
                new_shape.text = shape.text_frame.text
            else:
                if shape.shape_type == 13:  # 画像のShapeTypeは13
                    # 画像を複製
                    image_stream = shape.image.blob
                    new_slide.shapes.add_picture(
                        image_stream,
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )
                else:
                    # その他のシェイプ（例: グラフやテーブル）の処理が必要
                    pass

        return new_slide
    
    def insert_images_to_ppt(self):
        folder = self.folder_path.get()
        ppt_file = self.ppt_path.get()

        if not folder or not ppt_file:
            messagebox.showerror("エラー", "画像フォルダとPowerPointファイルを選択してください。")
            return

        # ユーザーが設定した画像サイズと位置を取得 (cmからEMUに変換)
        cm_to_emu = 360000  # 1cm = 360,000 EMU
        try:
            width = float(self.img_width.get()) * cm_to_emu
            height = float(self.img_height.get()) * cm_to_emu
            left = float(self.img_left.get()) * cm_to_emu
            top = float(self.img_top.get()) * cm_to_emu
        except ValueError:
            messagebox.showerror("エラー", "画像サイズや位置の入力に無効な値があります。")
            return

        # PowerPointファイルを開く
        presentation = Presentation(ppt_file)

        # 1ページ目のスライドレイアウトを複製
        first_slide_layout = presentation.slide_layouts[6]  # 1ページ目のスライドレイアウトを取得
        # first_slide_layout = presentation.Masters.get_Item(0)  # 1ページ目のスライドレイアウトを取得
        # 1ページ目のスライドを取得
        # first_slide = presentation.slides[0]
        
        # 画像の削除と新しい画像の挿入を行うスライドのインデックスを管理
        slides_with_images = {}  # スライド番号と画像オブジェクトのマッピング
        slide_index = 0

        for i, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # 画像のShapeTypeは13
                    slides_with_images[slide_index] = slide
                    slide_index += 1
                    break  # スライド内に1つ以上画像があれば、次のスライドへ

        # フォルダ内の画像を1枚ずつスライドに挿入または更新
        slide_index = 0
        for image_file in os.listdir(folder):
            if image_file.lower().endswith((".png", ".jpg", ".jpeg")):
                image_path = os.path.join(folder, image_file)

                # 画像を追加または更新するスライドを決定
                if slide_index in slides_with_images:
                    # 既存のスライドがある場合、画像を削除
                    slide = slides_with_images[slide_index]
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # 画像のShapeTypeは13
                            sp = shape
                            sp.element.getparent().remove(sp.element)  # 画像を削除
                            break
                else:
                    # 新しいスライドを追加
                    # slide_layout = presentation.slide_layouts[6]  # 空白スライド
                    # slide = presentation.slides.add_slide(slide_layout)
                    # 1ページ目のレイアウトを複製する
                    slide_layout = first_slide_layout
                    slide = presentation.slides.add_slide(slide_layout)
                    
                    # 新しいスライドを追加（1ページ目のスライドを複製）
                    # slide = self.clone_slide(first_slide, presentation)

                # 画像を挿入（ユーザーが指定した位置とサイズ）
                slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                slide_index += 1

        # PowerPointファイルを上書き保存
        presentation.save(ppt_file)
        messagebox.showinfo("完了", "画像がPowerPointに挿入されました。")

# GUIの実行
root = tk.Tk()
app = SlideMateApp(root)
root.mainloop()
